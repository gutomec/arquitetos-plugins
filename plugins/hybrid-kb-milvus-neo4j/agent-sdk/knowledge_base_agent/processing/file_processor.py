"""
File Processor - Document Processing Pipeline
==============================================
Processes uploaded files: extraction, chunking, embedding, indexing.
"""

import io
import hashlib
from typing import List, Optional
from uuid import UUID
import uuid
import logging
from datetime import datetime

from ..core.config import KBConfig
from ..core.models import DatabaseManager, KBFile, KBChunk, KnowledgeBase, FileStatus
from ..storage import StorageManager
from .embeddings import EmbeddingService
from .graph import GraphService

logger = logging.getLogger(__name__)


class FileProcessor:
    """Processes files through the complete ingestion pipeline."""

    def __init__(
        self,
        config: KBConfig,
        storage: StorageManager,
        embeddings: EmbeddingService,
        graph: GraphService,
        db: DatabaseManager
    ):
        self.config = config
        self.storage = storage
        self.embeddings = embeddings
        self.graph = graph
        self.db = db

    async def process_file(self, file_id: str) -> dict:
        """
        Process a file through the complete pipeline:
        1. Download from Minio
        2. Extract text
        3. Chunk text
        4. Generate embeddings
        5. Store in Milvus
        6. Create graph nodes in Neo4j
        7. Update PostgreSQL records
        """
        session = self.db.get_session()

        try:
            # Get file record
            file = session.query(KBFile).filter(KBFile.id == UUID(file_id)).first()
            if not file:
                return {"success": False, "error": "File not found"}

            kb = file.knowledge_base

            # Update status
            file.status = FileStatus.PROCESSING
            session.commit()

            # Download file
            content = self.storage.download_file(file.minio_object_key)

            # Handle ZIP files
            if file.file_type.value == "zip":
                return await self._process_zip(file, content, kb, session)

            # Extract text
            text = await self._extract_text(content, file.file_type.value, file.filename)

            if not text:
                file.status = FileStatus.FAILED
                file.error_message = "Could not extract text from file"
                session.commit()
                return {"success": False, "error": "Text extraction failed"}

            # Store preview
            file.text_preview = text[:500]

            # Chunk text
            chunks = self._chunk_text(text, kb.chunk_size, kb.chunk_overlap)

            # Create document node in Neo4j
            await self.graph.create_document_node(
                doc_id=str(file.id),
                kb_id=str(kb.id),
                filename=file.filename,
                metadata=file.metadata
            )

            # Process chunks
            vectors_to_insert = []
            chunk_records = []

            for i, chunk_text in enumerate(chunks):
                chunk_id = str(uuid.uuid4())
                content_hash = hashlib.sha256(chunk_text.encode()).hexdigest()

                # Generate embedding
                embedding = await self.embeddings.generate_embedding(
                    chunk_text,
                    kb.embedding_provider
                )

                # Prepare vector for Milvus
                milvus_id = f"{file_id}_{i}"
                vectors_to_insert.append({
                    "id": milvus_id,
                    "vector": embedding,
                    "text": chunk_text,
                    "file_id": str(file.id),
                    "kb_id": str(kb.id),
                    "chunk_index": i,
                    "metadata": {"filename": file.filename}
                })

                # Create chunk node in Neo4j
                await self.graph.create_chunk_node(
                    chunk_id=chunk_id,
                    doc_id=str(file.id),
                    content=chunk_text,
                    chunk_index=i,
                    milvus_id=milvus_id
                )

                # Create chunk record
                chunk_record = KBChunk(
                    id=UUID(chunk_id),
                    file_id=file.id,
                    chunk_index=i,
                    content=chunk_text,
                    content_hash=content_hash,
                    milvus_id=milvus_id,
                    neo4j_node_id=chunk_id
                )
                chunk_records.append(chunk_record)

            # Insert vectors into Milvus
            await self.embeddings.insert_vectors(kb.milvus_collection, vectors_to_insert)

            # Extract entities
            entity_count = await self.graph.extract_and_create_entities(
                str(file.id),
                text[:5000]  # Limit for performance
            )

            # Save chunk records
            session.add_all(chunk_records)

            # Update file record
            file.status = FileStatus.COMPLETED
            file.chunk_count = len(chunks)
            file.entity_count = entity_count
            file.processed_at = datetime.utcnow()

            # Update KB stats
            kb.file_count += 1
            kb.chunk_count += len(chunks)
            kb.total_size_bytes += file.size_bytes

            session.commit()

            logger.info(f"Processed file {file.filename}: {len(chunks)} chunks, {entity_count} entities")

            return {
                "success": True,
                "file_id": str(file.id),
                "chunks_created": len(chunks),
                "entities_extracted": entity_count
            }

        except Exception as e:
            logger.error(f"File processing failed: {e}")
            session.rollback()

            # Update file status
            try:
                file = session.query(KBFile).filter(KBFile.id == UUID(file_id)).first()
                if file:
                    file.status = FileStatus.FAILED
                    file.error_message = str(e)
                    session.commit()
            except:
                pass

            return {"success": False, "error": str(e)}

        finally:
            session.close()

    async def _process_zip(self, file: KBFile, content: bytes, kb: KnowledgeBase, session) -> dict:
        """Process a ZIP file by extracting and processing each file."""
        # Save ZIP to Minio first if not already there
        extracted = self.storage.extract_zip(file.minio_object_key)

        total_chunks = 0
        total_entities = 0
        processed_files = 0
        errors = []

        for filename, file_content, mime_type in extracted:
            try:
                # Determine file type
                ext = filename.lower().split('.')[-1] if '.' in filename else 'other'

                # Skip unsupported files
                if ext not in ['pdf', 'docx', 'xlsx', 'pptx', 'txt', 'md', 'csv', 'json', 'html']:
                    continue

                # Create child file record
                from ..core.models import FileType
                file_type = FileType(ext) if ext in [e.value for e in FileType] else FileType.OTHER

                child_file = KBFile(
                    knowledge_base_id=kb.id,
                    filename=filename,
                    original_filename=filename,
                    file_type=file_type,
                    mime_type=mime_type,
                    size_bytes=len(file_content),
                    status=FileStatus.PENDING,
                    parent_file_id=file.id,
                    is_from_zip=True
                )

                session.add(child_file)
                session.flush()

                # Upload to Minio
                object_key = f"{file.minio_object_key.rsplit('/', 1)[0]}/extracted/{child_file.id}/{filename}"
                self.storage.upload_file(object_key, file_content)
                child_file.minio_object_key = object_key

                session.commit()

                # Process the extracted file
                result = await self.process_file(str(child_file.id))

                if result.get("success"):
                    processed_files += 1
                    total_chunks += result.get("chunks_created", 0)
                    total_entities += result.get("entities_extracted", 0)
                else:
                    errors.append(f"{filename}: {result.get('error')}")

            except Exception as e:
                errors.append(f"{filename}: {str(e)}")

        # Update parent file status
        file.status = FileStatus.COMPLETED
        file.chunk_count = total_chunks
        file.entity_count = total_entities
        file.processed_at = datetime.utcnow()
        session.commit()

        return {
            "success": True,
            "file_id": str(file.id),
            "files_extracted": len(extracted),
            "files_processed": processed_files,
            "total_chunks": total_chunks,
            "total_entities": total_entities,
            "errors": errors if errors else None
        }

    async def _extract_text(self, content: bytes, file_type: str, filename: str) -> str:
        """Extract text from file content."""
        try:
            if file_type in ['txt', 'md', 'csv']:
                return content.decode('utf-8', errors='ignore')

            elif file_type == 'json':
                import json
                data = json.loads(content.decode('utf-8'))
                return json.dumps(data, indent=2)

            elif file_type == 'html':
                from html.parser import HTMLParser

                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []

                    def handle_data(self, data):
                        self.text.append(data.strip())

                parser = TextExtractor()
                parser.feed(content.decode('utf-8', errors='ignore'))
                return ' '.join(parser.text)

            elif file_type == 'pdf':
                try:
                    import pymupdf
                    doc = pymupdf.open(stream=content, filetype="pdf")
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    doc.close()
                    return text
                except ImportError:
                    logger.warning("pymupdf not installed, using fallback")
                    return f"[PDF content from {filename}]"

            elif file_type == 'docx':
                try:
                    from docx import Document
                    doc = Document(io.BytesIO(content))
                    return '\n'.join([para.text for para in doc.paragraphs])
                except ImportError:
                    logger.warning("python-docx not installed")
                    return f"[DOCX content from {filename}]"

            elif file_type == 'xlsx':
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
                    text = []
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            text.append(' | '.join(str(c) for c in row if c))
                    return '\n'.join(text)
                except ImportError:
                    logger.warning("openpyxl not installed")
                    return f"[XLSX content from {filename}]"

            elif file_type == 'pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return '\n'.join(text)
                except ImportError:
                    logger.warning("python-pptx not installed")
                    return f"[PPTX content from {filename}]"

            else:
                return content.decode('utf-8', errors='ignore')

        except Exception as e:
            logger.error(f"Text extraction error: {e}")
            return ""

    def _chunk_text(
        self,
        text: str,
        chunk_size: int = 512,
        overlap: int = 50
    ) -> List[str]:
        """Split text into chunks with overlap."""
        if not text:
            return []

        # Split by sentences
        import re
        sentences = re.split(r'(?<=[.!?])\s+', text)

        chunks = []
        current_chunk = ""

        for sentence in sentences:
            if len(current_chunk) + len(sentence) < chunk_size:
                current_chunk += sentence + " "
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                # Start new chunk with overlap
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else ""
                current_chunk = overlap_text + sentence + " "

        if current_chunk.strip():
            chunks.append(current_chunk.strip())

        return chunks if chunks else [text[:chunk_size]]
